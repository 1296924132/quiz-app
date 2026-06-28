from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色方案 ==========
BLUE_DARK = RGBColor(0x0D, 0x2B, 0x4E)      # 深蓝
BLUE_MID = RGBColor(0x1A, 0x4E, 0x8A)       # 中蓝
BLUE_LIGHT = RGBColor(0x3A, 0x7B, 0xD5)     # 亮蓝
ACCENT_CYAN = RGBColor(0x00, 0xBC, 0xD4)    # 青
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GOLD = RGBColor(0xFF, 0xB7, 0x00)

def add_bg(slide, color=BLUE_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, BLUE_DARK)

# 顶部装饰线
add_shape(slide1, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_CYAN)

# 左侧大色块装饰
add_shape(slide1, Inches(0), Inches(1.5), Inches(0.5), Inches(4.5), BLUE_MID)

# 主标题
add_text_box(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "科技兴警", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_text_box(slide1, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "智慧警务赋能平安中国建设", font_size=24, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.LEFT)

# 装饰分隔线
add_accent_line(slide1, Inches(1.5), Inches(3.8), Inches(2), GOLD)

# 副信息
add_text_box(slide1, Inches(1.5), Inches(4.2), Inches(6), Inches(0.5),
             "科技创新与警务实战深度融合", font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))

# 底部信息
add_text_box(slide1, Inches(1.5), Inches(6.5), Inches(6), Inches(0.4),
             "2026年6月", font_size=14, color=RGBColor(0x88, 0x99, 0xAA))

# 右侧圆形装饰
circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(2.5), Inches(2.5), Inches(2.5))
circle.fill.solid()
circle.fill.fore_color.rgb = BLUE_MID
circle.line.fill.background()
circle.fill.fore_color.brightness = 0.0

# 圆内文字
tf = circle.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "AI +\n警务"
p.font.size = Pt(20)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "微软雅黑"
p.alignment = PP_ALIGN.CENTER

# ============================================================
# 第2页：背景与意义
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# 顶部标题栏
add_shape(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BLUE_DARK)
add_text_box(slide2, Inches(0.8), Inches(0.2), Inches(10), Inches(0.6),
             "一、背景与意义", font_size=28, color=WHITE, bold=True)
add_accent_line(slide2, Inches(0.8), Inches(0.75), Inches(1.5), ACCENT_CYAN)

# 左半部分：政策背景
add_shape(slide2, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.05), BLUE_LIGHT)
add_text_box(slide2, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.5),
             "▎政策背景", font_size=20, color=BLUE_DARK, bold=True)

bg_items = [
    "🔹 公安部\"科技兴警\"三年行动计划（2023-2026）全面实施",
    "🔹 国家\"十四五\"规划将智慧警务列为重点方向",
    "🔹 数字中国建设整体布局规划推动政府治理数字化",
    "🔹 人工智能、大数据、物联网等技术日趋成熟",
]
add_bullet_list(slide2, Inches(0.6), Inches(2.1), Inches(5.8), Inches(3.5), bg_items, font_size=15)

# 右半部分：战略意义
add_shape(slide2, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.05), GOLD)
add_text_box(slide2, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.5),
             "▎战略意义", font_size=20, color=BLUE_DARK, bold=True)

importance_items = [
    "🔸 提升警务效能：科技手段替代传统人海战术",
    "🔸 增强预警能力：数据分析驱动风险精准预判",
    "🔸 优化资源配置：智能调度实现警力科学部署",
    "🔸 服务民生需求：\"互联网+警务\"便民利民",
]
add_bullet_list(slide2, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.5), importance_items, font_size=15)

# 底部总结框
add_shape(slide2, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0), LIGHT_GRAY)
add_text_box(slide2, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.7),
             "💡 核心定位：科技是公安工作现代化的重要引擎，是提升核心战斗力的关键支撑",
             font_size=15, color=BLUE_MID, bold=True)

# ============================================================
# 第3页：核心技术
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_shape(slide3, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BLUE_DARK)
add_text_box(slide3, Inches(0.8), Inches(0.2), Inches(10), Inches(0.6),
             "二、核心技术体系", font_size=28, color=WHITE, bold=True)
add_accent_line(slide3, Inches(0.8), Inches(0.75), Inches(1.5), ACCENT_CYAN)

# 4个核心技术卡片
techs = [
    ("人工智能", "🧠", [
        "视频图像智能分析",
        "语音语义识别",
        "行为异常检测",
        "案件智能研判",
    ]),
    ("大数据", "📊", [
        "多源数据融合分析",
        "犯罪规律挖掘",
        "人员关系图谱",
        "实时态势感知",
    ]),
    ("物联网", "🌐", [
        "智能感知设备网络",
        "移动警务终端",
        "电子围栏与定位",
        "物证智能管理",
    ]),
    ("云计算", "☁️", [
        "警务云平台建设",
        "弹性计算资源调度",
        "数据安全存储",
        "跨区域协同作战",
    ]),
]

card_width = Inches(2.8)
card_height = Inches(4.2)
gap = Inches(0.3)
start_x = Inches(0.6)
start_y = Inches(1.5)

for i, (title, icon, items) in enumerate(techs):
    x = start_x + i * (card_width + gap)
    # 卡片背景
    card = add_shape(slide3, x, start_y, card_width, card_height, WHITE)
    card.shadow.inherit = False
    # 卡片边框
    border = add_shape(slide3, x, start_y, card_width, Inches(0.06), BLUE_LIGHT)
    # 标题区域
    add_shape(slide3, x, start_y + Inches(0.06), card_width, Inches(0.7), LIGHT_GRAY)
    add_text_box(slide3, x + Inches(0.15), start_y + Inches(0.12), card_width - Inches(0.3), Inches(0.6),
                 f"{icon}  {title}", font_size=18, color=BLUE_DARK, bold=True)
    # 内容
    add_bullet_list(slide3, x + Inches(0.2), start_y + Inches(1.0), card_width - Inches(0.4), Inches(3.0),
                    ["• " + item for item in items], font_size=13, color=DARK_TEXT, spacing=Pt(10))

# ============================================================
# 第4页：应用场景
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

add_shape(slide4, Inches(0), Inches(0), Inches(13.333), Inches(1.0), BLUE_DARK)
add_text_box(slide4, Inches(0.8), Inches(0.2), Inches(10), Inches(0.6),
             "三、典型应用场景", font_size=28, color=WHITE, bold=True)
add_accent_line(slide4, Inches(0.8), Inches(0.75), Inches(1.5), ACCENT_CYAN)

scenarios = [
    ("🔍 智慧侦查", "AI辅助视频侦查\n人脸识别比对\n轨迹追踪分析\n大幅缩短破案周期", Inches(0.6), Inches(1.3)),
    ("🛡️ 智慧防控", "重点区域智能监控\n异常行为自动预警\n无人机自动巡防\n立体化治安防控体系", Inches(3.5), Inches(1.3)),
    ("🚔 智慧指挥", "可视化指挥调度\n警力资源一张图\n突发事件智能预案\n跨部门协同作战", Inches(6.4), Inches(1.3)),
    ("👮 智慧服务", "\"互联网+政务服务\"\n线上办事便民平台\n智能咨询机器人\n群众满意度提升", Inches(9.3), Inches(1.3)),
]

for title, desc, x, y in scenarios:
    # 场景卡片
    card = add_shape(slide4, x, y, Inches(2.7), Inches(3.5), WHITE)
    # 顶部色条
    add_shape(slide4, x, y, Inches(2.7), Inches(0.06), BLUE_LIGHT)
    # 标题
    add_text_box(slide4, x + Inches(0.15), y + Inches(0.2), Inches(2.4), Inches(0.5),
                 title, font_size=17, color=BLUE_DARK, bold=True)
    # 分割线
    add_accent_line(slide4, x + Inches(0.15), y + Inches(0.7), Inches(0.8), ACCENT_CYAN)
    # 描述
    add_text_box(slide4, x + Inches(0.2), y + Inches(0.9), Inches(2.3), Inches(2.4),
                 desc, font_size=13, color=DARK_TEXT)

# 底部总结
add_shape(slide4, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.8), BLUE_MID)
add_text_box(slide4, Inches(1.0), Inches(5.35), Inches(11.3), Inches(0.5),
             '科技赋能警务实战，实现"汗水警务"向"智慧警务"的跨越式转变',
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第5页：未来展望
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, BLUE_DARK)

add_shape(slide5, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
             "四、未来展望", font_size=32, color=WHITE, bold=True)
add_accent_line(slide5, Inches(0.8), Inches(1.05), Inches(1.5), GOLD)

# 三大方向
directions = [
    ("🤖 深化AI融合", "大模型技术深度赋能警务\n智能体自主研判与决策\n多模态AI全方位感知"),
    ("🔗 数据要素驱动", "数据要素市场化配置\n跨部门数据安全共享\n隐私计算保障数据安全"),
    ("🌍 全域智慧警务", "空地一体全域防控\n警企合作生态共建\n国际警务协作数字化"),
]

for i, (title, desc) in enumerate(directions):
    x = Inches(0.8) + i * Inches(4.1)
    # 数字标识
    num_shape = add_shape(slide5, x, Inches(1.6), Inches(0.5), Inches(0.5), GOLD)
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # 标题
    add_text_box(slide5, x + Inches(0.7), Inches(1.6), Inches(3.0), Inches(0.5),
                 title, font_size=20, color=GOLD, bold=True)
    # 分割线
    add_accent_line(slide5, x, Inches(2.3), Inches(1.2), ACCENT_CYAN)
    # 内容
    add_text_box(slide5, x, Inches(2.6), Inches(3.6), Inches(2.5),
                 desc, font_size=15, color=RGBColor(0xCC, 0xDD, 0xEE))

# 底部结束语
add_shape(slide5, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.06), ACCENT_CYAN)
add_text_box(slide5, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
             "科技兴警 智创未来  ·  让警务更智能 让社会更平安",
             font_size=18, color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)

# 页码
add_text_box(slide5, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
             "5 / 5", font_size=11, color=RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.RIGHT)

# ========== 给前4页加页码 ==========
for i in range(4):
    slide = list(prs.slides)[i]
    add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
                 f"{i + 1} / 5", font_size=11, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.RIGHT)

# ========== 保存 ==========
output_path = os.path.expanduser("~/科技兴警-智慧警务.pptx")
prs.save(output_path)
print(f"PPT已生成：{output_path}")
